"""Build final P2 presentation slides."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports" / "P2"
FIGURES = REPORTS / "figures"
OUT = REPORTS / "p2_final_slides.pptx"

INK = RGBColor(25, 34, 41)
MUTED = RGBColor(91, 101, 112)
TEAL = RGBColor(47, 111, 109)
RUST = RGBColor(196, 90, 58)
BLUE = RGBColor(79, 111, 168)
LINE = RGBColor(218, 225, 231)
BG = RGBColor(248, 250, 248)


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return slide


def _text(slide, x, y, w, h, text, size=22, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _title(slide, title, subtitle=None):
    _text(slide, 0.55, 0.28, 12.0, 0.48, title, size=24, bold=True)
    if subtitle:
        _text(slide, 0.57, 0.78, 11.6, 0.34, subtitle, size=11, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.12), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def _bullets(slide, x, y, w, h, items, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)
    return box


def _metric_card(slide, x, y, label, value, color):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(2.65), Inches(1.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    _text(slide, x + 0.15, y + 0.18, 2.35, 0.28, label, size=11, color=RGBColor(255, 255, 255), bold=True)
    _text(slide, x + 0.15, y + 0.53, 2.35, 0.42, value, size=24, color=RGBColor(255, 255, 255), bold=True)


def _picture(slide, path, x, y, w, h):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        _text(slide, x, y, w, 0.4, f"Missing figure: {path.name}", size=12, color=RUST)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _blank(prs)
    _text(slide, 0.7, 0.8, 11.8, 0.75, "P2 MLOps Final Results", size=34, bold=True)
    _text(slide, 0.72, 1.65, 11.5, 0.4, "PIU risk classification: tracking, search, registry, serving package, and visual evidence", size=15, color=MUTED)
    _metric_card(slide, 0.8, 2.65, "Baseline QWK", "0.3651", TEAL)
    _metric_card(slide, 3.7, 2.65, "Champion QWK", "0.3672", BLUE)
    _metric_card(slide, 6.6, 2.65, "Formal Trials", "10", RUST)
    _metric_card(slide, 9.5, 2.65, "Registry Aliases", "3", TEAL)
    _bullets(slide, 0.82, 4.3, 11.8, 1.4, [
        "Champion URI: models:/piu-risk@champion",
        "Formal study: p2-formal-mlops-20260612",
        "Main claim: complete MLOps loop; model lift is small and reported honestly.",
    ], size=18)

    slide = _blank(prs)
    _title(slide, "Pipeline Evidence", "From data split to MLflow registry and API-compatible model")
    _bullets(slide, 0.8, 1.55, 5.5, 4.8, [
        "Canonical 5-fold split and feature tables are reused from P1.",
        "Optuna explores model family, feature version, and hyperparameters.",
        "Each formal trial logs metrics to MLflow; final model is registered.",
        "baseline/candidate/champion aliases point to auditable model versions.",
        "Reports and figures are regenerated from local SQLite stores.",
    ], size=17)
    _picture(slide, FIGURES / "p2_mlflow_run_status.png", 6.7, 1.55, 5.6, 3.9)

    slide = _blank(prs)
    _title(slide, "Model Metrics", "Baseline vs selected champion")
    _picture(slide, FIGURES / "p2_metric_comparison.png", 0.85, 1.45, 7.0, 4.15)
    _bullets(slide, 8.1, 1.55, 4.4, 4.5, [
        "Baseline: Logistic Regression on feat_v1.",
        "Champion: Optuna-tuned Logistic Regression on feat_v1.",
        "QWK improves from 0.3651 to 0.3672.",
        "The gain is modest; the MLOps workflow is the main deliverable.",
    ], size=16)

    slide = _blank(prs)
    _title(slide, "Optuna Search Trace", "Trial objective and best-so-far curve")
    _picture(slide, FIGURES / "p2_optuna_trace.png", 0.75, 1.35, 7.2, 4.35)
    _bullets(slide, 8.15, 1.5, 4.4, 4.4, [
        "10 local formal trials were run.",
        "Best trial: logreg_v1, C=1.1159.",
        "MLP trials were fast but underperformed LR.",
        "feat_v2 did not improve this final champion.",
    ], size=16)

    slide = _blank(prs)
    _title(slide, "Champion Error Profile", "5-fold CV confusion matrix")
    _picture(slide, FIGURES / "p2_champion_confusion_matrix.png", 1.0, 1.3, 5.6, 4.9)
    _bullets(slide, 7.0, 1.55, 5.4, 4.2, [
        "The class imbalance remains visible.",
        "Most errors are between adjacent severity levels.",
        "This supports reporting QWK as the primary metric.",
        "Further model work should target minority classes and calibration.",
    ], size=16)

    slide = _blank(prs)
    _title(slide, "Deliverables", "What is ready for review")
    _bullets(slide, 0.85, 1.45, 11.8, 4.6, [
        "reports/P2/p2_final_report.md and p2_summary_report.md",
        "reports/P2/figures/*.png and *.svg",
        "reports/P2/p2_model_comparison.csv and p2_formal_optuna_trials.csv",
        "MLflow model versions 4/5/6 with baseline/candidate/champion aliases",
        "Champion can be loaded through models:/piu-risk@champion",
    ], size=18)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
