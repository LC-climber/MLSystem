"""
Build the P1 midterm presentation deck.

Requires:
  pip install python-pptx

Output:
  reports/P1/p1_midterm_slides.pptx
"""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports" / "P1"
FIGURES = REPORTS / "figures"
OUT = REPORTS / "p1_midterm_slides.pptx"

INK = RGBColor(23, 32, 42)
MUTED = RGBColor(95, 107, 122)
BLUE = RGBColor(76, 120, 168)
RED = RGBColor(228, 87, 86)
GREEN = RGBColor(84, 162, 75)
ORANGE = RGBColor(245, 133, 24)
LINE = RGBColor(217, 225, 234)
BG = RGBColor(247, 249, 251)


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return slide


def _text(slide, x, y, w, h, text, size=24, color=INK, bold=False,
          align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _title(slide, text, subtitle=None):
    _text(slide, 0.55, 0.28, 11.9, 0.45, text, size=24, color=INK, bold=True)
    if subtitle:
        _text(slide, 0.57, 0.78, 11.7, 0.32, subtitle, size=10.5, color=MUTED)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.08), Inches(12.1), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def _bullets(slide, x, y, w, h, items, size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        p.line_spacing = 1.1
    return box


def _callout(slide, x, y, w, h, text, color=ORANGE):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 248, 239)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = INK
    return shape


def _metric(slide, x, y, w, h, value, label, color=BLUE):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 251, 255)
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = MUTED
    return shape


def _picture_fit(slide, path, x, y, w, h):
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        width = w
        height = w / img_ratio
    else:
        height = h
        width = h * img_ratio
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(width), height=Inches(height))


def _footer(slide, page):
    _text(slide, 0.55, 7.12, 5.2, 0.22, "MLsystem P1 midterm | PIU multi-system comparison",
          size=8.5, color=MUTED)
    _text(slide, 12.25, 7.12, 0.55, 0.22, str(page), size=8.5, color=MUTED,
          align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = []

    s = _blank(prs)
    _text(s, 0.65, 0.85, 10.8, 0.7, "P1 中期汇报", size=34, bold=True)
    _text(s, 0.68, 1.55, 11.2, 0.45, "PIU 多系统算法对比:Spark 应该放在训练阶段,还是特征/ETL 阶段?", size=18, color=MUTED)
    _metric(s, 0.75, 2.45, 2.7, 1.1, "3 systems", "sklearn / Spark / PyTorch", BLUE)
    _metric(s, 3.75, 2.45, 2.7, 1.1, "2 algos", "LR + MLP", GREEN)
    _metric(s, 6.75, 2.45, 2.7, 1.1, "315M rows", "actigraphy feature stage", ORANGE)
    _metric(s, 9.75, 2.45, 2.7, 1.1, "A5/A6", "coverage + parallelism", RED)
    _callout(s, 0.75, 4.25, 11.75, 1.1,
             "核心结论:Spark 训练/单样本推理阶段不划算;特征/ETL 阶段也必须匹配正确算法和并行度。")
    slides.append(s)

    s = _blank(prs)
    _title(s, "1. 研究问题与公平协议", "同一任务,同一 fold,同一预处理;比较系统质量与系统成本")
    _bullets(s, 0.75, 1.45, 5.8, 4.4, [
        "任务:PIU 四分类,预测 sii=0/1/2/3",
        "系统:sklearn / Spark MLlib / PyTorch",
        "算法:Logistic Regression + MLP",
        "指标:Macro-F1 主指标,QWK / Balanced Accuracy / 系统指标为辅",
        "反泄漏:PCIAT-* 排除;imputer/scaler 只在训练 fold fit",
    ])
    _picture_fit(s, FIGURES / "p1_feature_lineage.png", 6.6, 1.35, 6.1, 4.4)
    slides.append(s)

    s = _blank(prs)
    _title(s, "2. 数据事实:为什么 Spark 有机会,也为什么 v2 难提升")
    _metric(s, 0.75, 1.35, 2.55, 1.0, "3960", "train feature rows", BLUE)
    _metric(s, 3.55, 1.35, 2.55, 1.0, "2736", "labeled CV rows", GREEN)
    _metric(s, 6.35, 1.35, 2.55, 1.0, "996", "actigraphy participants", ORANGE)
    _metric(s, 9.15, 1.35, 2.55, 1.0, "36.4%", "labeled act coverage", RED)
    _bullets(s, 0.85, 2.85, 5.55, 2.5, [
        "actigraphy 原始时序约 3.15 亿行,特征阶段确实有大数据量",
        "但有标签样本中只有 996/2736 有真实 actigraphy",
        "fold 4 的 actigraphy 子集没有 class 3 验证样本",
    ], size=16)
    _picture_fit(s, FIGURES / "p1_a5_coverage.png", 6.45, 2.55, 6.15, 3.6)
    slides.append(s)

    s = _blank(prs)
    _title(s, "3. Table 1:特征阶段 pandas vs Spark")
    _metric(s, 0.75, 1.4, 2.6, 1.05, "38.13s", "pandas streaming", BLUE)
    _metric(s, 3.65, 1.4, 2.6, 1.05, "0.65GB", "pandas RSS", BLUE)
    _metric(s, 6.55, 1.4, 2.6, 1.05, "114.01s", "Spark applyInPandas", RED)
    _metric(s, 9.45, 1.4, 2.6, 1.05, "13.25GB", "Spark RSS", RED)
    _bullets(s, 0.85, 3.05, 5.7, 2.6, [
        "naive pandas 全量读取会 OOM;按 id 分区流式聚合后稳定",
        "Spark 原生 exact percentile 路径 GC/OOM;改用 applyInPandas",
        "CPU/Spark 产物 max_abs_diff=1.14e-13,等价性通过",
    ], size=16)
    _picture_fit(s, FIGURES / "p1_feature_lineage.png", 6.4, 2.85, 6.2, 2.9)
    _callout(s, 0.85, 6.05, 11.7, 0.55, "结论:Spark 的合理位置在特征/ETL,但实现必须算法感知。", ORANGE)
    slides.append(s)

    s = _blank(prs)
    _title(s, "4. Table 2:模型质量没有因 feat_v2 稳定提升")
    _picture_fit(s, FIGURES / "p1_table2_metric_bars.png", 0.65, 1.15, 12.1, 4.15)
    _bullets(s, 0.9, 5.55, 11.4, 1.0, [
        "LR 三系统 Macro-F1 在 v2/all 均略降;MLP 只有 Spark MLP 明显上升",
        "v2/all 的主协议结论:actigraphy 特征没有转化为跨系统稳定收益",
    ], size=15.5)
    slides.append(s)

    s = _blank(prs)
    _title(s, "5. 系统成本:Spark 不适合小表格训练和单样本推理")
    _picture_fit(s, FIGURES / "p1_system_costs.png", 0.85, 1.25, 7.4, 4.65)
    _bullets(s, 8.45, 1.45, 3.85, 3.75, [
        "Spark LR v2 训练约 25.3s/fold",
        "Spark 单样本推理约 156ms/job",
        "sklearn/PyTorch 推理在几十到几百微秒量级",
        "小表格在线推理阶段,Spark 调度成本压倒计算本身",
    ], size=16)
    _callout(s, 8.45, 5.55, 3.85, 0.75, "系统结论比模型分数更关键。", RED)
    slides.append(s)

    s = _blank(prs)
    _title(s, "6. A5:覆盖率解释 feat_v2 为什么没有稳定收益")
    _picture_fit(s, FIGURES / "p1_a5_coverage.png", 0.7, 1.25, 7.05, 4.6)
    _bullets(s, 8.0, 1.35, 4.45, 4.45, [
        "全体 feature 行覆盖:996/3960 = 25.2%",
        "有标签 CV 样本覆盖:996/2736 = 36.4%",
        "actigraphy 子集 fold 4 class 3 = 0",
        "所以 v2/actigraphy 只能作补充,不能替代 Table 2 主表",
    ], size=16)
    slides.append(s)

    s = _blank(prs)
    _title(s, "7. A6:Spark 并行度不是越高越快")
    _picture_fit(s, FIGURES / "p1_a6_spark_parallelism.png", 0.9, 1.25, 6.55, 4.75)
    _metric(s, 8.0, 1.45, 3.5, 0.95, "local[4]", "fastest setting", BLUE)
    _metric(s, 8.0, 2.75, 3.5, 0.95, "115.48s", "wall time", ORANGE)
    _metric(s, 8.0, 4.05, 3.5, 0.95, "12.42GB", "peak RSS", RED)
    _callout(s, 7.75, 5.55, 4.5, 0.8, "瓶颈更像 shuffle/序列化/Python worker 调度,不是 CPU core 数不足。", ORANGE)
    slides.append(s)

    s = _blank(prs)
    _title(s, "8. 代表性混淆矩阵:类别 3 是指标不稳定来源")
    _picture_fit(s, FIGURES / "p1_confusion_sklearn_lr.png", 0.9, 1.15, 11.5, 5.2)
    _callout(s, 1.05, 6.35, 11.1, 0.55, "class 3 极少,且常被预测为 class 2;Macro-F1 对这种尾类很敏感。", RED)
    slides.append(s)

    s = _blank(prs)
    _title(s, "9. 汇报主线:Spark 价值在哪里")
    _bullets(s, 0.9, 1.35, 5.9, 4.75, [
        "训练阶段:Spark 在小表格上慢,推理尤其不适合单样本在线场景",
        "特征阶段:actigraphy 有真实大数据量,但 exact 聚合要算法感知",
        "并行度:local[4] 快于 local[8]/local[20],必须实测拐点",
        "质量:feat_v2 没有稳定提升,主要受覆盖率与类别稀疏限制",
    ], size=17)
    _picture_fit(s, FIGURES / "p1_feature_lineage.png", 6.75, 1.55, 5.8, 3.9)
    slides.append(s)

    s = _blank(prs)
    _title(s, "10. 结论与下一步")
    _callout(s, 0.85, 1.3, 11.5, 0.9,
             "P1 核心结论已成立:系统选择不能只看模型分数,必须把训练/推理/特征工程成本放进同一条证据链。")
    _bullets(s, 0.95, 2.65, 5.55, 3.25, [
        "已完成:Table 1, Table 2, A5, A6,核心可视化,动画 HTML",
        "可选补充:A1-A4 轻量消融作为附录",
        "下一阶段:P2 MLOps,模型注册/服务化/监控/双发布",
    ], size=17)
    _bullets(s, 7.0, 2.65, 5.05, 3.25, [
        "交付物:00_docs/P1_MIDTERM_REPORT.md",
        "交付物:reports/P1/p1_midterm_slides.pptx",
        "交付物:reports/P1/p1_midterm_explainer.html",
        "交付物:reports/P1/figures/",
    ], size=15.5, color=MUTED)
    slides.append(s)

    s = _blank(prs)
    _title(s, "Appendix A:答辩 Q&A 速记")
    _bullets(s, 0.85, 1.3, 5.8, 4.9, [
        "为什么 Spark 没赢? 小表格训练/单样本推理阶段调度开销大。",
        "为什么仍说 Spark 价值在 ETL? actigraphy 有约 3.15 亿行,特征阶段才有大数据量。",
        "为什么 feat_v2 没提升? 有标签样本 actigraphy 覆盖率只有 36.4%,fold 4 子集无 class 3。",
        "为什么不用 actigraphy 子集替代主表? 样本协议改变,类别更稀疏,只能作补充分析。",
    ], size=15.5)
    _bullets(s, 7.0, 1.3, 5.4, 4.9, [
        "A6 为什么越并行越慢? 瓶颈更像 shuffle/序列化/Python worker 调度和尾部任务。",
        "公平性怎么保证? 同 feature、fold、seed、预处理;imputer/scaler 只在训练 fold fit。",
        "P2 怎么衔接? 选择稳定低延迟模型,做注册、服务化、监控和发布。",
        "完整回答见:00_docs/P1_MIDTERM_QA.md",
    ], size=15.5, color=MUTED)
    slides.append(s)

    s = _blank(prs)
    _title(s, "Appendix B:复现命令与材料路径")
    _bullets(s, 0.85, 1.25, 5.8, 4.95, [
        "Table 1: python -m src.experiments.run_p1_feature_stage --mlflow",
        "Table 2: python -m src.experiments.run_p1_systemwise --feature v2 --mlflow",
        "A5: python -m src.experiments.run_p1_ablation_a5_coverage",
        "A6: python -m src.experiments.run_p1_spark_parallelism --mlflow",
        "Figures: python -m src.experiments.run_p1_visualizations",
        "PPT: python scripts/build_p1_midterm_ppt.py",
    ], size=13.5)
    _bullets(s, 7.0, 1.25, 5.45, 4.95, [
        "报告:00_docs/P1_MIDTERM_REPORT.md",
        "讲稿:00_docs/P1_MIDTERM_TALK_TRACK.md",
        "Q&A:00_docs/P1_MIDTERM_QA.md",
        "PPT:reports/P1/p1_midterm_slides.pptx",
        "HTML:reports/P1/p1_midterm_explainer.html",
        "图件:reports/P1/figures/",
    ], size=14.5, color=MUTED)
    slides.append(s)

    for i, slide in enumerate(slides, 1):
        _footer(slide, i)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT} ({len(slides)} slides)")


if __name__ == "__main__":
    build()
